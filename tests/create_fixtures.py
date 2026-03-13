"""Generate binary test fixture files (PDF, XLSX, PPTX, DOCX).

Run once: python tests/create_fixtures.py
Requires: pymupdf, openpyxl, python-pptx, python-docx

This script is temporary — delete after Phase 4A testing is complete.
"""
from pathlib import Path

FIXTURES = Path(__file__).parent / "fixtures"


def create_pdf():
    import pymupdf
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "テスト用PDFドキュメントです。")
    doc.save(str(FIXTURES / "sample.pdf"))
    doc.close()


def create_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "テスト"
    ws["A1"] = "項目"
    ws["B1"] = "値"
    ws["A2"] = "テスト項目1"
    ws["B2"] = 100
    wb.save(str(FIXTURES / "sample.xlsx"))


def create_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "テストスライド"
    slide.placeholders[1].text = "テスト用プレゼンテーションです。"
    prs.save(str(FIXTURES / "sample.pptx"))


def create_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("テストドキュメント", level=1)
    doc.add_paragraph("テスト用Wordドキュメントです。")
    doc.save(str(FIXTURES / "sample.docx"))


if __name__ == "__main__":
    create_pdf()
    create_xlsx()
    create_pptx()
    create_docx()
    print("Fixtures created.")
